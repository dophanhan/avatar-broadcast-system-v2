# -*- coding: utf-8 -*-
"""
AI 数字人播报系统 V3.0 - 阿里云函数计算后端 (修复致命 Bug 版)
集成：阿里云 OSS + 豆包大模型 (火山引擎)

修复清单：
1. ✅ 补充 handle_project_save 和 handle_project_delete 中的 OSS 环境变量
2. ✅ 修正阿里云 IMM 签名算法（使用标准 percent_encode）
3. ✅ 修正 GetRangeRequest 缺少 inclusive_start_primary_key 和 exclusive_end_primary_key
"""

# 添加 python 目录到 PATH（阿里云 FC 要求）
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'python'))

import json
import hashlib
import hmac
import base64
from datetime import datetime, timezone
from urllib.parse import quote
import os
import uuid
import time
import requests
import oss2
from pptx import Presentation
from io import BytesIO
from tablestore import OTSClient, OTSError, SingleColumnCondition, ComparatorType, Condition, Row, ReturnType, INF_MIN, INF_MAX, RowExistenceExpectation

# 统一的 CORS 响应头
CORS_HEADERS = {
    'Access-Control-Allow-Origin': '*',
    'Access-Control-Allow-Methods': 'POST, GET, OPTIONS, PUT, DELETE',
    'Access-Control-Allow-Headers': 'Content-Type, X-Feishu-OpenId, Authorization, Accept',
    'Access-Control-Max-Age': '86400',
    'Access-Control-Expose-Headers': 'Date, x-fc-request-id',
    'Content-Type': 'application/json;charset=utf-8'
}

def handler(event, context):
    """阿里云函数计算 Python 3.10 入口（兼容 FC 3.0 和 FC 2.0）"""
    if isinstance(event, bytes):
        event = event.decode('utf-8')
    if isinstance(event, str):
        event = json.loads(event)
    
    # ✅ 修复路由解析：兼容 FC 3.0 和 FC 2.0
    # FC 3.0 使用 rawPath 和 requestContext.http.method
    # FC 2.0 使用 path 和 httpMethod/method
    path = event.get('rawPath', event.get('path', '/'))
    
    request_context = event.get('requestContext', {})
    http_context = request_context.get('http', {})
    method = http_context.get('method', event.get('httpMethod', event.get('method', 'POST')))
    
    # OPTIONS 跨域预检请求
    if method == 'OPTIONS':
        return {'statusCode': 200, 'headers': CORS_HEADERS, 'body': ''}
    
    # 路由分发
    if path == '/ai/generate-turning-nodes' and method == 'POST':
        return handle_generate_turning_nodes(event)
    elif path == '/ai/generate-broadcast-text' and method == 'POST':
        return handle_generate_broadcast_text(event)
    elif path == '/project/save' and method == 'POST':
        return handle_project_save(event)
    elif path == '/project/list' and method == 'GET':
        return handle_project_list(event)
    elif path == '/project/detail' and method == 'GET':
        return handle_project_detail(event)
    elif path == '/project/delete' and method == 'POST':
        return handle_project_delete(event)
    elif path == '/ppt/upload' and method == 'POST':
        return handle_ppt_upload(event)
    elif path == '/doc/parse' and method == 'POST':
        return handle_doc_parse(event)
    else:
        return handle_avatar_auth(event)

def make_response(code, data, error_msg=None):
    """统一的响应构造器"""
    if code == 0:
        body = {"code": 0, "data": data}
    else:
        body = {"code": code, "msg": error_msg or "Unknown error"}
    return {
        'statusCode': 200 if code == 0 else 500,
        'headers': CORS_HEADERS,
        'body': json.dumps(body, ensure_ascii=False)
    }

def handle_avatar_auth(event):
    """讯飞数字人鉴权接口"""
    try:
        API_KEY = os.environ.get('AVATAR_API_KEY', '')
        API_SECRET = os.environ.get('AVATAR_API_SECRET', '')
        APP_ID = os.environ.get('AVATAR_APP_ID', '')
        SCENE_ID = os.environ.get('AVATAR_SCENE_ID', '')
        WS_URL = "wss://avatar.cn-huadong-1.xf-yun.com/v1/interact"
        
        date = datetime.utcnow().strftime('%a, %d %b %Y %H:%M:%S GMT')
        signature_origin = f"host: avatar.cn-huadong-1.xf-yun.com\ndate: {date}\nGET /v1/interact HTTP/1.1"
        signature_sha = hmac.new(API_SECRET.encode(), signature_origin.encode(), digestmod=hashlib.sha256).digest()
        signature = base64.b64encode(signature_sha).decode()
        authorization_origin = f'api_key="{API_KEY}", algorithm="hmac-sha256", headers="host date request-line", signature="{signature}"'
        authorization = base64.b64encode(authorization_origin.encode()).decode()
        signed_url = f"{WS_URL}?authorization={quote(authorization)}&date={quote(date)}&host=avatar.cn-huadong-1.xf-yun.com"
        
        body = event.get('body', '{}')
        if isinstance(body, bytes):
            body = body.decode('utf-8')
        data = json.loads(body) if body else {}
        
        response_data = {
            "appId": APP_ID,
            "sceneId": SCENE_ID,
            "signedUrl": signed_url,
            "avatarId": data.get("avatarId"),
            "vcn": data.get("vcn"),
            "speakQueue": []
        }
        return make_response(0, response_data)
    except Exception as e:
        return make_response(-1, None, str(e))

def handle_generate_turning_nodes(event):
    """智能翻页节点生成 (豆包大模型版 - 增强 Prompt)"""
    try:
        body = json.loads(event.get('body', '{}'))
        # 👇 兼容新旧两种参数格式
        prompt = body.get("prompt", "")
        sentences = body.get("sentences", body.get("speak_sentences", []))
        page_count = body.get("page_count", 0)
        ppt_pages = body.get("ppt_pages", [])
        
        # 如果是新格式（带 prompt），直接使用
        if prompt:
            user_prompt = prompt
        elif sentences and ppt_pages:
            # 旧格式：使用原有逻辑
            user_prompt = f"PPT 页面数据：{json.dumps(ppt_pages, ensure_ascii=False)}\n\n播报文案句子列表：{json.dumps(sentences, ensure_ascii=False)}"
        else:
            raise Exception('缺少播报句子或 PPT 数据')

        DOUBAO_API_KEY = os.environ.get('DOUBAO_API_KEY', '')
        DOUBAO_MODEL_ID = os.environ.get('DOUBAO_MODEL_ID', '')  # ⚠️ 豆包推理接入点 ID: ep-m-20260224164411-wv277
        
        if not DOUBAO_API_KEY or not DOUBAO_MODEL_ID:
            raise Exception('未配置豆包 API Key 或推理接入点 ID')

        system_prompt = """你是一个专业的视频导播和 PPT 播报节奏分析专家。
你的任务是根据用户提供的【PPT 每一页的内容】和【播报文案句子列表】，计算出在播报哪几句话时，应该显示哪一页 PPT。

规则：
1. 句子索引从 0 开始。
2. 每一页 PPT 必须连续对应一段句子（即有明确的起始句和结束句）。
3. 所有句子必须被完全分配，不能有遗漏，也不能重复分配。
4. 直接返回纯 JSON 对象，不要任何 Markdown 格式（如 ```json），不要任何解释文字。

输出格式：
{
    "page_turning_nodes": [
        {"page_num": 1, "start_sentence_index": 0, "end_sentence_index": 2},
        {"page_num": 2, "start_sentence_index": 3, "end_sentence_index": 5}
    ]
}

注意：只返回 JSON 对象本身，不要任何其他内容。"""

        headers = {
            "Authorization": f"Bearer {DOUBAO_API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": DOUBAO_MODEL_ID,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            "temperature": 0.1
        }
        
        url = "https://ark.cn-beijing.volces.com/api/v3/chat/completions"
        res = requests.post(url, headers=headers, json=payload, timeout=300)  # 10 分钟超时
        res_data = res.json()
        
        if res.status_code != 200 or "choices" not in res_data:
            raise Exception(f"豆包 API 调用失败：{res_data}")
            
        ai_response_text = res_data["choices"][0]["message"]["content"]
        
        # 🛡️ 终极防御：暴力剔除大模型可能带有的 Markdown 格式包裹
        ai_response_text = ai_response_text.strip()
        if ai_response_text.startswith("```json"):
            ai_response_text = ai_response_text[7:]
        elif ai_response_text.startswith("```"):
            ai_response_text = ai_response_text[3:]
        if ai_response_text.endswith("```"):
            ai_response_text = ai_response_text[:-3]
        ai_response_text = ai_response_text.strip()
        
        # 解析 JSON 结果
        nodes_data = json.loads(ai_response_text)
        
        return make_response(0, {"page_turning_nodes": nodes_data.get("page_turning_nodes", [])})
        
    except Exception as e:
        print(f"[ERROR] AI 节点生成失败：{str(e)}")
        return make_response(-1, None, str(e))

# ========== 阿里云 IMM 标准签名算法 (修复 Bug 2) ==========
def percent_encode(encode_str):
    """阿里云专属的 URL 编码方法"""
    import urllib.parse
    res = urllib.parse.quote(str(encode_str), safe='')
    res = res.replace('+', '%20').replace('*', '%2A').replace('%7E', '~')
    return res

def compute_imm_signature(params, access_key_secret):
    """计算阿里云 IMM 签名"""
    sorted_params_list = sorted(params.items(), key=lambda x: x[0])
    canonicalized_query_string = '&'.join([f'{percent_encode(k)}={percent_encode(v)}' for k, v in sorted_params_list])
    string_to_sign = 'POST&%2F&' + percent_encode(canonicalized_query_string)
    
    signature = hmac.new(
        (access_key_secret + "&").encode('utf-8'),
        string_to_sign.encode('utf-8'),
        hashlib.sha1
    ).digest()
    return base64.b64encode(signature).decode('utf-8'), canonicalized_query_string

def handle_ppt_upload(event):
    """本地 PPT 文件上传解析 + 阿里云 OSS 存储 + IMM 文档转换"""
    try:
        # FC 3.0 兼容性处理：body 可能是 base64 编码的字符串
        body_raw = event.get('body', '{}')
        is_base64_encoded = event.get('isBase64Encoded', False)
        
        # 如果是 base64 编码，先解码
        if is_base64_encoded and isinstance(body_raw, str):
            body_str = base64.b64decode(body_raw).decode('utf-8')
        else:
            body_str = body_raw if isinstance(body_raw, str) else body_raw.decode('utf-8')
        
        # 解析 JSON
        data = json.loads(body_str)
        file_base64 = data.get('file', '')
        filename = data.get('filename', 'presentation.pptx')
        
        if ',' in file_base64:
            file_data = base64.b64decode(file_base64.split(',')[1])
        else:
            file_data = base64.b64decode(file_base64)
            
        ppt_id = str(uuid.uuid4())
        
        # --- 1. 上传源文件到阿里云 OSS ---
        OSS_ACCESS_KEY = os.environ.get('OSS_ACCESS_KEY', '')
        OSS_SECRET_KEY = os.environ.get('OSS_SECRET_KEY', '')
        OSS_ENDPOINT = os.environ.get('OSS_ENDPOINT', '')
        OSS_BUCKET_NAME = os.environ.get('OSS_BUCKET_NAME', '')
        OSS_REGION = os.environ.get('OSS_REGION', 'cn-hangzhou')
        
        if OSS_ACCESS_KEY and OSS_BUCKET_NAME:
            auth = oss2.Auth(OSS_ACCESS_KEY, OSS_SECRET_KEY)
            bucket = oss2.Bucket(auth, OSS_ENDPOINT, OSS_BUCKET_NAME)
            
            oss_object_key = f"ppt_source/{ppt_id}/{filename}"
            bucket.put_object(oss_object_key, file_data)
            print(f"[SUCCESS] 源文件已上传至 OSS: {oss_object_key}")
            
            # --- 2. 调用阿里云 IMM 进行 PPT 转图片 (修复 Bug 2: 标准签名算法) ---
            try:
                IMM_ENDPOINT = f"https://imm.{OSS_REGION}.aliyuncs.com"
                src_uri = f"oss://{OSS_BUCKET_NAME}/{oss_object_key}"
                dst_uri = f"oss://{OSS_BUCKET_NAME}/ppt_images/{ppt_id}/"
                
                # CreateOfficeConversionTask 请求参数 (2020-09-30 规范)
                params = {
                    "Action": "CreateOfficeConversionTask",
                    "Version": "2020-09-30",
                    "ProjectName": "default",
                    "SourceURI": src_uri,
                    
                    # ✅ 使用 TargetURIPrefix，阿里云自动加上 1.jpg, 2.jpg...
                    "TargetURIPrefix": dst_uri,
                    
                    "TargetType": "jpg",
                    "PageRanges": "1-",
                    "Scale": "100",
                    "Timestamp": datetime.now(timezone.utc).strftime('%Y-%m-%dT%H:%M:%SZ'),
                    "AccessKeyId": OSS_ACCESS_KEY,
                    "SignatureMethod": "HMAC-SHA1",
                    "SignatureVersion": "1.0",
                    "SignatureNonce": str(uuid.uuid4()),
                    "Format": "JSON"
                }
                
                # 使用标准签名算法
                signature_b64, canonicalized_query_string = compute_imm_signature(params, OSS_SECRET_KEY)
                final_data = canonicalized_query_string + f"&Signature={percent_encode(signature_b64)}"
                
                imm_headers = {
                    "Content-Type": "application/x-www-form-urlencoded"
                }
                
                # 发送 IMM 请求
                imm_res = requests.post(IMM_ENDPOINT, headers=imm_headers, data=final_data, timeout=10)
                print(f"[IMM] 响应状态码：{imm_res.status_code}")
                print(f"[IMM] 原始响应：{imm_res.text}")  # 💡 打印原始响应，便于调试
                
                # 检查响应内容
                if not imm_res.text or imm_res.text.strip() == '':
                    print(f"[WARN] IMM 返回空响应")
                    imm_result = {}
                else:
                    try:
                        imm_result = imm_res.json()
                        print(f"[IMM] 解析成功：{imm_result}")
                    except json.JSONDecodeError as e:
                        print(f"[WARN] IMM 响应解析失败：{e}")
                        print(f"[WARN] 原始响应内容：{imm_res.text[:500]}")
                        imm_result = {}
                
                if imm_res.status_code == 200 and "TaskId" in imm_result:
                    task_id = imm_result["TaskId"]
                    print(f"[SUCCESS] IMM 转换任务已提交：{task_id}")
                    
                    # 轮询等待转换完成（最多等待 120 秒 = 2 分钟）
                    for _ in range(120):
                        # 轮询查询任务状态
                        check_params = {
                            # ✅ 2020-09-30 版统一的查询接口
                            "Action": "GetTask",
                            "Version": "2020-09-30",
                            "ProjectName": "default",
                            
                            # ✅ 新版接口强制要求提供任务类型
                            "TaskType": "OfficeConversion",
                            
                            "TaskId": task_id,
                            "Timestamp": datetime.now(timezone.utc).strftime('%Y-%m-%dT%H:%M:%SZ'),
                            "AccessKeyId": OSS_ACCESS_KEY,
                            "SignatureMethod": "HMAC-SHA1",
                            "SignatureVersion": "1.0",
                            "SignatureNonce": str(uuid.uuid4()),
                            "Format": "JSON"
                        }
                        
                        # 使用标准签名算法
                        check_sig_b64, check_canonicalized = compute_imm_signature(check_params, OSS_SECRET_KEY)
                        check_final_data = check_canonicalized + f"&Signature={percent_encode(check_sig_b64)}"
                        
                        check_res = requests.post(IMM_ENDPOINT, headers=imm_headers, data=check_final_data, timeout=10)
                        check_result = check_res.json()
                        
                        status = check_result.get("Status")
                        print(f"[IMM] 当前轮询状态：{status}")  # 💡 打印状态
                        
                        # ✅ 同时兼容旧版的 Finished 和新版的 Succeeded
                        if status in ["Finished", "Succeeded"]:
                            print("[SUCCESS] IMM 转换完成")
                            break
                        elif status == "Failed":
                            print("[WARN] IMM 转换失败")
                            break
                        
                        time.sleep(1)  # 睡 1 秒继续问
                            
                else:
                    print(f"[WARN] IMM 任务提交失败：{imm_result}")
                    
            except Exception as imm_error:
                print(f"[WARN] IMM 调用异常：{str(imm_error)}")
            
        # --- 3. 内存解析 PPT 文本内容 ---
        ppt_stream = BytesIO(file_data)
        prs = Presentation(ppt_stream)
        
        pages = []
        for i, slide in enumerate(prs.slides):
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            page_text = " ".join(text_content) if text_content else f"第{i+1}页"
            title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"第{i+1}页"
            
            if i == 0: page_type = "cover"
            elif i == len(prs.slides) - 1: page_type = "ending"
            else: page_type = "content"
            
            image_url = ""
            if OSS_ACCESS_KEY and OSS_BUCKET_NAME:
                image_url = f"https://{OSS_BUCKET_NAME}.{OSS_ENDPOINT}/ppt_images/{ppt_id}/{i+1}.jpg"

            pages.append({
                "page_num": i + 1,
                "page_title": title,
                "page_text": page_text[:500],
                "page_image_url": image_url, 
                "page_type": page_type
            })
        
        response_data = {
            "ppt_id": ppt_id,
            "ppt_title": filename.replace('.pptx', '').replace('.ppt', ''),
            "total_pages": len(pages),
            "pages": pages
        }
        
        return make_response(0, response_data)
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"[ERROR] PPT 解析失败：{str(e)}")
        print(f"[ERROR] Traceback: {error_trace}")
        return make_response(-1, None, f"PPT 解析失败：{str(e)}")

def handle_project_save(event):
    """项目保存 - 阿里云表格存储 + OSS"""
    try:
        body = json.loads(event.get('body', '{}'))
        project_id = body.get("project_id", str(uuid.uuid4()))
        project_name = body.get("project_name", "未命名项目")
        project_desc = body.get("project_desc", "")
        avatar_config = body.get("avatar_config", {})
        speak_text = body.get("speak_text", "")
        speak_sentences = body.get("speak_sentences", [])
        ppt_data = body.get("ppt_data", {})
        page_turning_rules = body.get("page_turning_rules", [])
        video_file = body.get("video_file")
        
        # 获取表格存储配置
        OTS_ENDPOINT = os.environ.get('OTS_ENDPOINT', '')
        OTS_INSTANCE_NAME = os.environ.get('OTS_INSTANCE_NAME', '')
        OTS_ACCESS_KEY = os.environ.get('OTS_ACCESS_KEY', '')
        OTS_SECRET_KEY = os.environ.get('OTS_SECRET_KEY', '')
        
        # 👇 ========== 新增：OTS 协议头绝对防御机制 ========== 👇
        if OTS_ENDPOINT and not OTS_ENDPOINT.startswith('http'):
            print(f"[INFO] 自动修复 OTS_ENDPOINT 协议头：{OTS_ENDPOINT}")
            OTS_ENDPOINT = f"https://{OTS_ENDPOINT}"
        # 👆 =================================================== 👆
        
        # 【修复 Bug 1】补充 OSS 环境变量
        OSS_ACCESS_KEY = os.environ.get('OSS_ACCESS_KEY', '')
        OSS_SECRET_KEY = os.environ.get('OSS_SECRET_KEY', '')
        OSS_ENDPOINT = os.environ.get('OSS_ENDPOINT', '')
        OSS_BUCKET_NAME = os.environ.get('OSS_BUCKET_NAME', '')
        
        # 保存视频到 OSS（如果有）
        video_url = ""
        if video_file and OSS_ACCESS_KEY and OSS_BUCKET_NAME:
            try:
                auth = oss2.Auth(OSS_ACCESS_KEY, OSS_SECRET_KEY)
                bucket = oss2.Bucket(auth, OSS_ENDPOINT, OSS_BUCKET_NAME)
                
                video_key = f"project_videos/{project_id}/video.webm"
                if ',' in video_file:
                    video_data = base64.b64decode(video_file.split(',')[1])
                else:
                    video_data = base64.b64decode(video_file)
                
                bucket.put_object(video_key, video_data)
                video_url = f"https://{OSS_BUCKET_NAME}.{OSS_ENDPOINT}/{video_key}"
                print(f"[SUCCESS] 视频已上传至 OSS: {video_key}")
            except Exception as video_error:
                print(f"[WARN] 视频上传失败：{str(video_error)}")
        
        # 保存到表格存储
        if OTS_ENDPOINT and OTS_INSTANCE_NAME:
            try:
                ots_client = OTSClient(OTS_ACCESS_KEY, OTS_SECRET_KEY, OTS_ENDPOINT, OTS_INSTANCE_NAME)
                
                project_table = 'avatar_projects'
                timestamp = int(time.time() * 1000)
                
                project_data = {
                    'project_id': project_id,
                    'project_name': project_name,
                    'project_desc': project_desc,
                    'avatar_config': json.dumps(avatar_config, ensure_ascii=False),
                    'speak_text': speak_text,
                    'speak_sentences': json.dumps(speak_sentences, ensure_ascii=False),
                    'ppt_data': json.dumps(ppt_data, ensure_ascii=False),
                    'page_turning_rules': json.dumps(page_turning_rules, ensure_ascii=False),
                    'video_url': video_url,
                    'video_duration': ppt_data.get('totalPages', 0) * 30,
                    'create_time': datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S'),
                    'update_time': datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S'),
                    'timestamp': timestamp
                }
                
                from tablestore import Condition, ReturnType, Row, RowExistenceExpectation
                
                # 创建主键（列表格式）
                pk = [('project_id', project_id)]
                
                # 创建条件
                cond = Condition(RowExistenceExpectation.IGNORE)
                
                # 创建 Row 对象
                row = Row(pk, project_data)
                
                # 调用 put_row
                ots_client.put_row(project_table, row, cond, ReturnType.NONE)
                print(f"[SUCCESS] 项目已保存至表格存储：{project_id}")
                
            except OTSError as ots_e:
                print(f"[WARN] 表格存储操作失败：{str(ots_e)}")
            except Exception as ots_e:
                print(f"[WARN] 表格存储异常：{str(ots_e)}")
        
        return make_response(0, {
            "project_id": project_id,
            "video_url": video_url
        })
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"[ERROR] 项目保存失败：{str(e)}")
        print(f"[ERROR] Traceback: {error_trace}")
        return make_response(-1, None, f"项目保存失败：{str(e)}")

def handle_project_list(event):
    """项目列表查询 - 表格存储（支持 GET 和 POST）"""
    try:
        # 兼容 GET 和 POST 请求
        if event.get('method') == 'POST' or event.get('httpMethod') == 'POST':
            body_raw = event.get('body', '{}')
            if isinstance(body_raw, bytes):
                body_raw = body_raw.decode('utf-8')
            body = json.loads(body_raw) if body_raw else {}
        else:
            body = {}
        
        OTS_ENDPOINT = os.environ.get('OTS_ENDPOINT', '')
        OTS_INSTANCE_NAME = os.environ.get('OTS_INSTANCE_NAME', '')
        OTS_ACCESS_KEY = os.environ.get('OTS_ACCESS_KEY', '')
        OTS_SECRET_KEY = os.environ.get('OTS_SECRET_KEY', '')
        
        projects = []
        
        if OTS_ENDPOINT and OTS_INSTANCE_NAME:
            try:
                ots_client = OTSClient(OTS_ACCESS_KEY, OTS_SECRET_KEY, OTS_ENDPOINT, OTS_INSTANCE_NAME)
                # 【修复 Bug 3】引入 INF_MIN 和 INF_MAX
                from tablestore import GetRangeRequest, Direction, INF_MIN, INF_MAX
                
                # 设定主键的扫描范围（从无限小到无限大，代表全表扫描）
                inclusive_start_primary_key = [('project_id', INF_MIN)]
                exclusive_end_primary_key = [('project_id', INF_MAX)]
                
                get_range_request = GetRangeRequest(
                    table_name='avatar_projects',
                    direction=Direction.FORWARD,
                    inclusive_start_primary_key=inclusive_start_primary_key,
                    exclusive_end_primary_key=exclusive_end_primary_key,
                    limit=100
                )
                
                response = ots_client.get_range(get_range_request)
                
                for row in response.rows:
                    project = {}
                    for key, value in row.columns.items():
                        if key in ['avatar_config', 'speak_sentences', 'ppt_data', 'page_turning_rules']:
                            try:
                                project[key] = json.loads(value)
                            except:
                                project[key] = value
                        else:
                            project[key] = value
                    projects.append(project)
                
                projects.sort(key=lambda x: x.get('timestamp', 0), reverse=True)
                print(f"[SUCCESS] 查询到 {len(projects)} 个项目")
                
            except Exception as ots_e:
                print(f"[WARN] 表格存储查询失败：{str(ots_e)}")
        
        return make_response(0, {
            "total": len(projects),
            "list": projects
        })
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"[ERROR] 项目列表查询失败：{str(e)}")
        print(f"[ERROR] Traceback: {error_trace}")
        return make_response(0, {"total": 0, "list": []})

def handle_project_detail(event):
    """项目详情查询 - 表格存储"""
    try:
        project_id = event.get('queryParameters', {}).get('project_id', '')
        
        if not project_id:
            raise Exception('缺少 project_id 参数')
        
        OTS_ENDPOINT = os.environ.get('OTS_ENDPOINT', '')
        OTS_INSTANCE_NAME = os.environ.get('OTS_INSTANCE_NAME', '')
        OTS_ACCESS_KEY = os.environ.get('OTS_ACCESS_KEY', '')
        OTS_SECRET_KEY = os.environ.get('OTS_SECRET_KEY', '')
        
        project_data = {}
        
        if OTS_ENDPOINT and OTS_INSTANCE_NAME:
            try:
                ots_client = OTSClient(OTS_ACCESS_KEY, OTS_SECRET_KEY, OTS_ENDPOINT, OTS_INSTANCE_NAME)
                
                # 创建主键（列表格式）
                pk = [('project_id', project_id)]
                
                response = ots_client.get_row('avatar_projects', pk)
                
                for key, value in response.row.columns.items():
                    if key in ['avatar_config', 'speak_sentences', 'ppt_data', 'page_turning_rules']:
                        try:
                            project_data[key] = json.loads(value)
                        except:
                            project_data[key] = value
                    else:
                        project_data[key] = value
                
                print(f"[SUCCESS] 项目详情已加载：{project_id}")
                
            except Exception as ots_e:
                print(f"[WARN] 表格存储查询失败：{str(ots_e)}")
        
        return make_response(0, project_data)
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"[ERROR] 项目详情查询失败：{str(e)}")
        print(f"[ERROR] Traceback: {error_trace}")
        return make_response(-1, None, f"项目详情查询失败：{str(e)}")

def handle_project_delete(event):
    """项目删除 - 表格存储 + OSS"""
    try:
        body = json.loads(event.get('body', '{}'))
        project_id = body.get("project_id", "")
        
        if not project_id:
            raise Exception('缺少 project_id 参数')
        
        OTS_ENDPOINT = os.environ.get('OTS_ENDPOINT', '')
        OTS_INSTANCE_NAME = os.environ.get('OTS_INSTANCE_NAME', '')
        OTS_ACCESS_KEY = os.environ.get('OTS_ACCESS_KEY', '')
        OTS_SECRET_KEY = os.environ.get('OTS_SECRET_KEY', '')
        
        # 【修复 Bug 1】补充 OSS 环境变量
        OSS_ACCESS_KEY = os.environ.get('OSS_ACCESS_KEY', '')
        OSS_SECRET_KEY = os.environ.get('OSS_SECRET_KEY', '')
        OSS_ENDPOINT = os.environ.get('OSS_ENDPOINT', '')
        OSS_BUCKET_NAME = os.environ.get('OSS_BUCKET_NAME', '')
        
        # 删除表格存储记录
        if OTS_ENDPOINT and OTS_INSTANCE_NAME:
            try:
                ots_client = OTSClient(OTS_ACCESS_KEY, OTS_SECRET_KEY, OTS_ENDPOINT, OTS_INSTANCE_NAME)
                
                # 创建主键（列表格式）
                pk = [('project_id', project_id)]
                
                ots_client.delete_row('avatar_projects', pk)
                print(f"[SUCCESS] 项目已从表格存储删除：{project_id}")
                
            except Exception as ots_e:
                print(f"[WARN] 表格存储删除失败：{str(ots_e)}")
        
        # 删除 OSS 视频文件
        if OSS_ACCESS_KEY and OSS_BUCKET_NAME:
            try:
                auth = oss2.Auth(OSS_ACCESS_KEY, OSS_SECRET_KEY)
                bucket = oss2.Bucket(auth, OSS_ENDPOINT, OSS_BUCKET_NAME)
                
                video_key = f"project_videos/{project_id}/video.webm"
                bucket.delete_object(video_key)
                print(f"[SUCCESS] 视频已从 OSS 删除：{video_key}")
                
            except Exception as oss_e:
                print(f"[WARN] OSS 删除失败：{str(oss_e)}")
        
        return make_response(0, {})
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"[ERROR] 项目删除失败：{str(e)}")
        print(f"[ERROR] Traceback: {error_trace}")
        return make_response(-1, None, f"项目删除失败：{str(e)}")

def handle_doc_parse(event):
    """文档提取纯文本 (支持 PDF, DOCX, TXT)"""
    try:
        import base64
        import io
        import json
        
        # ⚠️ 关键修复：兼容阿里云 FC API 网关的 Base64 编码封装
        body_raw = event.get('body', '{}')
        is_base64_encoded = event.get('isBase64Encoded', False)
        
        if is_base64_encoded and isinstance(body_raw, str):
            body_str = base64.b64decode(body_raw).decode('utf-8')
        else:
            body_str = body_raw if isinstance(body_raw, str) else body_raw.decode('utf-8')
            
        body = json.loads(body_str)
        file_base64 = body.get('file', '')
        filename = body.get('filename', '').lower()
        
        if not file_base64:
            raise Exception("未接收到文件数据")

        # 处理前端传来的 Base64 前缀 (如 data:application/pdf;base64,...)
        if ',' in file_base64:
            file_data = base64.b64decode(file_base64.split(',')[1])
        else:
            file_data = base64.b64decode(file_base64)
            
        extracted_text = ""
        
        # 1. 解析 TXT
        if filename.endswith('.txt'):
            try:
                extracted_text = file_data.decode('utf-8')
            except UnicodeDecodeError:
                extracted_text = file_data.decode('gbk', errors='ignore')
                
        # 2. 解析 DOCX
        elif filename.endswith('.docx'):
            try:
                import docx
                doc = docx.Document(io.BytesIO(file_data))
                # 完美保留段落回车
                extracted_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            except ImportError:
                raise Exception("后端缺少 python-docx 依赖库，请在 requirements.txt 中添加")
                
        # 3. 解析 PDF
        elif filename.endswith('.pdf'):
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(io.BytesIO(file_data))
                pages_text = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        pages_text.append(text.strip())
                # 用双回车分隔不同页面，保证排版清晰
                extracted_text = "\n\n".join(pages_text)
            except ImportError:
                raise Exception("后端缺少 PyPDF2 依赖库，请在 requirements.txt 中添加")
                
        elif filename.endswith('.doc'):
            raise Exception("暂不支持老版本的 .doc 格式，请另存为 .docx 后上传")
        else:
            raise Exception(f"不支持的文件格式：{filename}")

        return make_response(0, {"parsed_text": extracted_text.strip()})
        
    except Exception as e:
        import traceback
        print(f"[ERROR] 文档解析失败：{str(e)}")
        print(traceback.format_exc())
        return make_response(-1, None, f"解析失败：{str(e)}")

def handle_generate_broadcast_text(event):
    """AI 生成播报文本 - 豆包大模型"""
    try:
        import json
        import requests
        
        body = json.loads(event.get('body', '{}'))
        ppt_pages = body.get('ppt_pages', [])
        ppt_title = body.get('ppt_title', '演示文稿')
        total_pages = body.get('total_pages', 0)
        
        if not ppt_pages or total_pages == 0:
            raise Exception('缺少 PPT 内容数据')
        
        # 获取豆包配置
        DOUBAO_API_KEY = os.environ.get('DOUBAO_API_KEY', '')
        DOUBAO_MODEL_ID = os.environ.get('DOUBAO_MODEL_ID', '')
        
        if not DOUBAO_API_KEY or not DOUBAO_MODEL_ID:
            raise Exception('未配置豆包 API Key 或 Model ID')
        
        # 构建 PPT 内容字符串
        ppt_content_formatted = "\n\n".join([
            f"第{page['page_num']}页 - {page.get('page_title', '无标题')}:\n{page.get('page_text', '无内容')}"
            for page in ppt_pages
        ])
        
        # 系统提示词（最终版）
        system_prompt = """你是一位顶级的编辑和播报文案撰写专家。
你的任务是：根据 PPT 核心信息，将其转化为完全适配数字人播报的口语化讲解稿。

请严格遵守以下【四项核心准则】：

<规则 1：数字人专属合规与适配>
- 标点限制：仅允许使用中文句号、逗号、顿号、冒号。禁止使用问号、感叹号、省略号、破折号、引号、书名号等。将所有疑问句、感叹句改写为语气平稳的陈述句。
- 字符与数字：所有数值、年份、百分比统一使用阿拉伯数字。绝对禁止出现任何英文字母或单词（遇到英文必须转化为标准中文意译，如 PPT->幻灯片，AI->人工智能，CEO->首席执行官）。
- 读音安全：全程规避生僻字、多音字、易读错字，替换为同义常用字。调整同音字引发的歧义表述。
- 内容安全：严禁包含政治、宗教、医疗、金融理财及个人隐私等敏感信息。严禁使用绝对化广告违规用语（如最好、第一、极品、百分百）。严禁包含网址、二维码、电话等引流信息。

<规则 2：播报节奏与文本结构>
- 结构框架：必须遵循"开场问候与主题引入 -> 核心内容串讲（逻辑连贯，自然过渡） -> 结尾致谢"的完整结构。
- 气息控制：单句长度严格控制在 15-30 字之间，确保数字人播报时自然断句，无长句卡顿。
- 语调风格：采用正式、平稳、亲切的口语化表达，将书面语转化为易于听懂的讲解语言。杜绝网络烂梗和方言。

<规则 3：内容转化策略>
- 信息提炼：覆盖 PPT 核心论点和关键数据，无需逐字复述。根据 PPT 信息密度自动调整文案长度，无遗漏、无冗余。

<规则 4：最高优先级输出格式（极其重要）>
- 必须是纯文本。绝对禁止输出任何 Markdown 标记（如 **加粗**、# 标题）、HTML 标签、序号（如 1. 2. 3.）、列表、前言后语、解释性备注。
- 除了最终的播报正文，不要说任何多余的话。

请直接输出符合要求的纯文本播报文案："""
        
        # 用户提示词
        user_prompt = f"""【PPT 标题】{ppt_title}

【PPT 页面内容】
{ppt_content_formatted}

请开始生成："""
        
        # 调用豆包 API
        headers = {
            "Authorization": f"Bearer {DOUBAO_API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": DOUBAO_MODEL_ID,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            "temperature": 0.7
        }
        
        url = "https://ark.cn-beijing.volces.com/api/v3/chat/completions"
        res = requests.post(url, headers=headers, json=payload, timeout=500)  # 500 秒超时，避免长文本生成失败
        res_data = res.json()
        
        if res.status_code != 200 or "choices" not in res_data:
            raise Exception(f"豆包 API 调用失败：{res_data}")
        
        broadcast_text = res_data["choices"][0]["message"]["content"]
        
        # 后处理：清理 Markdown 标记
        broadcast_text = broadcast_text.strip()
        if broadcast_text.startswith('```'):
            lines = broadcast_text.split('\n')
            broadcast_text = '\n'.join(lines[1:-1]) if len(lines) > 2 else broadcast_text
        
        return {
            'statusCode': 200,
            'headers': {'Access-Control-Allow-Origin': '*', 'Content-Type': 'application/json'},
            'body': json.dumps({
                'code': 0,
                'data': {
                    'broadcast_text': broadcast_text
                }
            }, ensure_ascii=False)
        }
        
    except Exception as e:
        print(f"[ERROR] AI 播报文本生成失败：{str(e)}")
        return {
            'statusCode': 500,
            'headers': {'Access-Control-Allow-Origin': '*', 'Content-Type': 'application/json'},
            'body': json.dumps({
                'code': -1,
                'msg': str(e)
            }, ensure_ascii=False)
        }
